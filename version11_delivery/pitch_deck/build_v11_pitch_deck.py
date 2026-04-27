from __future__ import annotations

import json
import math
import zipfile
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
ANALYTICS = ROOT / "version11_bert" / "outputs" / "analytics"
EVAL = ROOT / "version11_bert" / "eval_outputs"
PROCESSED = ROOT / "version11_bert" / "processed"
OUT_DIR = ROOT / "version11_delivery" / "pitch_deck"
OUT_PPTX = OUT_DIR / "dataleague_v11_final_pitch_deck.pptx"
OUT_PLAN = OUT_DIR / "V11_SLIDE_FLOW.md"

W, H = 13.333, 7.5

COLORS = {
    "bg": "F7F4EE",
    "ink": "111827",
    "muted": "5B6472",
    "line": "D8D1C4",
    "risk": "E6483D",
    "risk_dark": "A92B25",
    "trust": "0E7C66",
    "blue": "2D5BFF",
    "amber": "F2B84B",
    "panel": "FFFFFF",
    "soft": "EAE3D7",
    "dark": "172033",
}


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip("#")
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:], 16))


def set_bg(slide, color: str = COLORS["bg"]) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 20,
    color: str = COLORS["ink"],
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Aptos",
    valign=MSO_ANCHOR.TOP,
) -> object:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.62, 0.38, 8.8, 0.55, size=28, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.64, 0.92, 9.4, 0.38, size=12, color=COLORS["muted"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.64), Inches(1.32), Inches(1.3), Inches(0.045))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(COLORS["risk"])
    line.line.fill.background()


def add_footer(slide, idx: int) -> None:
    add_text(slide, "DataLeague Final | V11", 0.64, 7.12, 2.4, 0.2, size=8, color="7C8796")
    add_text(slide, f"{idx:02d}", 12.25, 7.1, 0.45, 0.22, size=9, color="7C8796", align=PP_ALIGN.RIGHT)


def add_chip(slide, text: str, x: float, y: float, w: float, fill: str, color: str = "FFFFFF") -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = rgb(color)


def add_metric(slide, value: str, label: str, x: float, y: float, w: float, accent: str) -> None:
    add_text(slide, value, x, y, w, 0.58, size=31, bold=True, color=accent)
    add_text(slide, label, x, y + 0.58, w, 0.36, size=10, color=COLORS["muted"])


def add_rule(slide, x: float, y: float, w: float, color: str = COLORS["line"]) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.018))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(color)
    line.line.fill.background()


def add_panel(slide, x: float, y: float, w: float, h: float, fill: str = COLORS["panel"]) -> object:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb(fill)
    panel.line.color.rgb = rgb(COLORS["line"])
    panel.line.width = Pt(0.7)
    return panel


def risk_color(value: float) -> str:
    if value >= 0.50:
        return COLORS["risk"]
    if value >= 0.30:
        return "F07C59"
    if value >= 0.18:
        return COLORS["amber"]
    if value >= 0.10:
        return "9AC9A8"
    return "DDE8DD"


def pct(v: float) -> str:
    return f"{v * 100:.1f}%"


def load_data() -> dict:
    summary = json.loads((ANALYTICS / "analytics_summary.json").read_text(encoding="utf-8"))
    metrics = json.loads((PROCESSED / "metrics.json").read_text(encoding="utf-8"))
    external_raw = json.loads((EVAL / "text_only_qwen_validation_2950_dedup_v11_metrics.json").read_text(encoding="utf-8"))
    external = external_raw.get("overall", external_raw)
    theme = pd.read_csv(ANALYTICS / "theme_risk.csv").fillna({"primary_theme": "Unknown"})
    pl = pd.read_csv(ANALYTICS / "platform_language_risk.csv")
    author = pd.read_csv(ANALYTICS / "author_risk.csv")
    trend = pd.read_csv(ANALYTICS / "daily_risk_trend.csv")
    dup = pd.read_csv(ANALYTICS / "high_risk_duplicate_clusters.csv")
    return {
        "summary": summary,
        "metrics": metrics,
        "external": external,
        "theme": theme,
        "pl": pl,
        "author": author,
        "trend": trend,
        "dup": dup,
    }


def style_chart(chart, font_size: int = 9) -> None:
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    for axis_name in ("category_axis", "value_axis"):
        try:
            axis = getattr(chart, axis_name)
            axis.tick_labels.font.size = Pt(font_size)
            axis.tick_labels.font.color.rgb = rgb(COLORS["muted"])
            axis.format.line.color.rgb = rgb("C7BFB1")
        except Exception:
            pass
    try:
        chart.value_axis.major_gridlines.format.line.color.rgb = rgb("E4DDD2")
    except Exception:
        pass


def build_deck() -> None:
    data = load_data()
    summary = data["summary"]
    metrics = data["metrics"]
    external = data["external"]
    theme = data["theme"]
    pl = data["pl"]
    author = data["author"]
    trend = data["trend"]
    dup = data["dup"]

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["dark"])
    add_text(slide, "V11", 0.72, 0.55, 2.0, 0.55, size=22, bold=True, color=COLORS["amber"])
    add_text(slide, "Sosyal Medya\nManipülasyon Tespit Sistemi", 0.72, 1.28, 8.2, 1.55, size=38, bold=True, color="FFFFFF")
    add_text(
        slide,
        "Text-first XLM-RoBERTa modeli + 5M içerik skorlama + manipülasyon haritası + canlı jüri tahmini",
        0.78,
        3.15,
        7.2,
        0.56,
        size=16,
        color="D7E0EC",
    )
    add_metric(slide, "5.0M", "skorlanan içerik", 0.8, 4.52, 2.0, COLORS["amber"])
    add_metric(slide, "0.943", "external ROC-AUC", 3.05, 4.52, 2.2, "FFFFFF")
    add_metric(slide, "89.9%", "M recall", 5.55, 4.52, 1.8, COLORS["risk"])
    for i, (label, val, color) in enumerate(
        [
            ("Platform", 0.78, COLORS["blue"]),
            ("Dil", 0.60, COLORS["risk"]),
            ("Tema", 0.42, COLORS["amber"]),
            ("Author", 0.70, COLORS["trust"]),
        ]
    ):
        x = 9.0 + i * 0.75
        h = 3.35 * val
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(5.75 - h), Inches(0.42), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(color)
        bar.line.fill.background()
        add_text(slide, label, x - 0.1, 5.92, 0.65, 0.23, size=8, color="C9D3E2", align=PP_ALIGN.CENTER)
    add_text(slide, "DataLeague Final", 0.78, 6.82, 2.2, 0.25, size=10, color="A9B5C5")
    add_text(slide, "Canlı demo hazır", 10.5, 6.82, 1.7, 0.25, size=10, bold=True, color=COLORS["amber"])

    # 2. Required outputs
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "İstenen üç çıktı aynı sistemde birleşiyor", "Model, harita ve canlı tahmin aynı V11 bundle’ını kullanıyor.")
    pillars = [
        ("01", "Organiklik skoru", "Her içerik için 0-1 arası manipülatif skor ve organiklik skoru.", COLORS["trust"]),
        ("02", "Manipülasyon haritası", "Platform, dil, tema, zaman, author ve duplicate cluster kırılımları.", COLORS["risk"]),
        ("03", "Canlı inference", "Jüri metni anlık işlenir; karar, skor, review flag ve gerekçe döner.", COLORS["blue"]),
    ]
    for i, (num, head, body, color) in enumerate(pillars):
        x = 0.78 + i * 4.15
        add_text(slide, num, x, 1.8, 0.75, 0.45, size=28, bold=True, color=color)
        add_text(slide, head, x, 2.38, 3.2, 0.4, size=20, bold=True)
        add_rule(slide, x, 2.92, 2.3, color)
        add_text(slide, body, x, 3.2, 3.35, 1.1, size=15, color=COLORS["muted"])
    add_text(slide, "Karar modeli text-only; metadata dashboard ve ağ analitiği için kullanıldı.", 1.1, 5.55, 10.8, 0.45, size=22, bold=True, color=COLORS["dark"], align=PP_ALIGN.CENTER)
    add_chip(slide, "original_text", 4.35, 6.18, 1.65, COLORS["dark"])
    add_chip(slide, "XLM-RoBERTa", 6.2, 6.18, 1.75, COLORS["risk"])
    add_chip(slide, "score + reason", 8.15, 6.18, 1.75, COLORS["trust"])
    add_footer(slide, 2)

    # 3. Architecture
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Mimari: canlı karar sade, analitik katman zengin", "Jüri unseen metin verdiğinde model yalnız metne bakar; 5M skorlar haritada bağlama oturur.")
    steps = [
        ("60k manuel M/O", "group split\nleakage kontrol"),
        ("XLM-RoBERTa", "CUDA eğitimi\nclass-weighted CE"),
        ("Eşik politikası", "M recall ≥ 0.80\nreview bandı"),
        ("5M scoring", "row-group shard\nresume + combine"),
        ("Dashboard", "platform/dil/tema\nuser/zaman/ağ"),
    ]
    for i, (head, body) in enumerate(steps):
        x = 0.65 + i * 2.5
        add_panel(slide, x, 2.12, 2.0, 1.12, "FFFFFF")
        add_text(slide, head, x + 0.18, 2.32, 1.65, 0.28, size=13, bold=True)
        add_text(slide, body, x + 0.18, 2.68, 1.65, 0.42, size=9, color=COLORS["muted"])
        if i < len(steps) - 1:
            add_text(slide, "→", x + 2.1, 2.45, 0.32, 0.3, size=22, bold=True, color=COLORS["risk"])
    add_text(slide, "Model kararı", 1.0, 4.15, 2.0, 0.26, size=13, bold=True, color=COLORS["risk"])
    add_text(slide, "original_text → manipulative_score / organic_score → Organik / Manipülatif", 1.0, 4.55, 7.1, 0.32, size=20, bold=True)
    add_text(slide, "Analitik karar desteği", 1.0, 5.25, 2.2, 0.26, size=13, bold=True, color=COLORS["trust"])
    add_text(slide, "Skorlar; platform, dil, tema, tarih, author_hash ve duplicate kümeleriyle gruplanır.", 1.0, 5.65, 9.1, 0.32, size=18, color=COLORS["muted"])
    add_footer(slide, 3)

    # 4. Model validation
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Model: manipülatifi kaçırmamaya ayarlı", "Eşik validation üzerinde M recall hedefiyle seçildi; canlı sonuç review flag ile güven bandı verir.")
    hold = metrics["bert"]["holdout"]
    ext = external
    chart_data = CategoryChartData()
    chart_data.categories = ["ROC-AUC", "PR-AUC", "M Recall", "M F1"]
    chart_data.add_series("Holdout", [hold["roc_auc"], hold["pr_auc"], hold["recall_M"], hold["f1_M"]])
    chart_data.add_series("External", [ext["roc_auc"], ext["pr_auc"], ext["recall_M"], ext["f1_M"]])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.75), Inches(1.75), Inches(6.35), Inches(4.35), chart_data).chart
    style_chart(chart, 9)
    chart.value_axis.maximum_scale = 1.0
    chart.value_axis.minimum_scale = 0.0
    chart.plots[0].series[0].format.fill.solid()
    chart.plots[0].series[0].format.fill.fore_color.rgb = rgb(COLORS["blue"])
    chart.plots[0].series[1].format.fill.solid()
    chart.plots[0].series[1].format.fill.fore_color.rgb = rgb(COLORS["risk"])
    add_metric(slide, "0.943", "External ROC-AUC", 7.65, 1.82, 2.1, COLORS["risk"])
    add_metric(slide, "0.892", "External PR-AUC", 10.05, 1.82, 2.0, COLORS["blue"])
    add_metric(slide, "89.9%", "External M recall", 7.65, 3.45, 2.2, COLORS["trust"])
    add_metric(slide, "78.2%", "External M F1", 10.05, 3.45, 2.0, COLORS["amber"])
    add_text(slide, "Confusion matrix external validation", 7.7, 5.05, 4.0, 0.25, size=12, bold=True)
    cm = ext["confusion_matrix"]["matrix"]
    x0, y0, cw, ch = 7.7, 5.42, 1.1, 0.42
    labels = [["", "Pred O", "Pred M"], ["Actual O", str(cm[0][0]), str(cm[0][1])], ["Actual M", str(cm[1][0]), str(cm[1][1])]]
    for r in range(3):
        for c in range(3):
            fill = "FFFFFF"
            if r == 2 and c == 2:
                fill = "DDEFE8"
            elif r == 1 and c == 2:
                fill = "F9D3CB"
            elif r == 2 and c == 1:
                fill = "F7E5B4"
            elif r == 0 or c == 0:
                fill = "EAE3D7"
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0 + c * cw), Inches(y0 + r * ch), Inches(cw), Inches(ch))
            rect.fill.solid()
            rect.fill.fore_color.rgb = rgb(fill)
            rect.line.color.rgb = rgb("FFFFFF")
            add_text(slide, labels[r][c], x0 + c * cw + 0.04, y0 + r * ch + 0.1, cw - 0.08, 0.2, size=9, bold=(r == 0 or c == 0), align=PP_ALIGN.CENTER)
    add_footer(slide, 4)

    # 5. Full scoring
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "5M veri üzerinde skorlama tamamlandı", "Her satıra skor, karar ve review flag basıldı; çıktı dashboard ve demo ile aynı eşiği kullanıyor.")
    counts = summary["decision_counts"]
    chart_data = ChartData()
    chart_data.categories = ["Organic", "Manipulative"]
    chart_data.add_series("decision", [counts["Organic"], counts["Manipulative"]])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.9), Inches(1.7), Inches(4.0), Inches(4.0), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_DATA_LABEL_POSITION.BEST_FIT
    chart.plots[0].series[0].points[0].format.fill.solid()
    chart.plots[0].series[0].points[0].format.fill.fore_color.rgb = rgb(COLORS["trust"])
    chart.plots[0].series[0].points[1].format.fill.solid()
    chart.plots[0].series[0].points[1].format.fill.fore_color.rgb = rgb(COLORS["risk"])
    add_metric(slide, f"{summary['rows_scored']:,}".replace(",", "."), "skorlanan içerik", 5.55, 1.9, 3.1, COLORS["dark"])
    add_metric(slide, pct(summary["decision_rates"]["manipulative_rate"]), "manipülatif karar oranı", 5.55, 3.38, 3.1, COLORS["risk"])
    add_metric(slide, pct(summary["decision_rates"]["review_rate"]), "review bandındaki içerik", 8.75, 3.38, 2.8, COLORS["amber"])
    add_text(slide, "Çıktı şeması", 5.55, 5.1, 1.4, 0.25, size=13, bold=True)
    add_text(slide, "row_id, original_text, language, url, author_hash, date, primary_theme, manipulative_score, organic_score, decision, review_flag", 5.55, 5.48, 6.4, 0.6, size=13, color=COLORS["muted"])
    add_footer(slide, 5)

    # 6. Platform language heatmap
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Manipülasyon haritası: risk platform ve dile göre ayrışıyor", "Yüksek hacimli kombinasyonlar filtrelendi; renk ortalama manipülatif skoru gösterir.")
    rows = ["x.com", "www.youtube.com", "www.reddit.com", "reddit.com", "bsky.app"]
    cols = ["en", "ar", "zh", "hi", "tr", "de", "ru"]
    pivot = pl.pivot_table(index="url", columns="language", values="mean_manipulative_score", aggfunc="mean")
    count_pivot = pl.pivot_table(index="url", columns="language", values="rows", aggfunc="sum")
    x0, y0 = 2.3, 1.85
    cell_w, cell_h = 1.05, 0.54
    add_text(slide, "Platform", 0.72, y0 + 0.12, 1.25, 0.2, size=10, bold=True, color=COLORS["muted"])
    for j, col in enumerate(cols):
        add_text(slide, col, x0 + j * cell_w, y0 - 0.35, cell_w, 0.22, size=10, bold=True, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    for i, row_name in enumerate(rows):
        add_text(slide, row_name, 0.72, y0 + i * cell_h + 0.14, 1.5, 0.18, size=10, bold=True)
        for j, col in enumerate(cols):
            val = float(pivot.loc[row_name, col]) if row_name in pivot.index and col in pivot.columns and not math.isnan(pivot.loc[row_name, col]) else 0.0
            rows_n = int(count_pivot.loc[row_name, col]) if row_name in count_pivot.index and col in count_pivot.columns and not math.isnan(count_pivot.loc[row_name, col]) else 0
            fill = risk_color(val) if rows_n >= 500 else "EFEAE2"
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0 + j * cell_w), Inches(y0 + i * cell_h), Inches(cell_w - 0.04), Inches(cell_h - 0.04))
            rect.fill.solid()
            rect.fill.fore_color.rgb = rgb(fill)
            rect.line.color.rgb = rgb("FFFFFF")
            label = f"{val:.2f}" if rows_n >= 500 else "-"
            add_text(slide, label, x0 + j * cell_w + 0.05, y0 + i * cell_h + 0.17, cell_w - 0.12, 0.16, size=8, bold=val >= 0.3, align=PP_ALIGN.CENTER, color=("FFFFFF" if val >= 0.5 and rows_n >= 500 else COLORS["ink"]))
    add_text(slide, "Öne çıkan bulgular", 9.95, 1.9, 2.1, 0.25, size=13, bold=True, color=COLORS["risk"])
    notes = [
        "x.com/ar: 68k içerik, ortalama skor 0.602",
        "YouTube/en: 110k içerik, M oranı 38.2%",
        "Reddit kırılımları geniş hacimli ama düşük riskli",
        "Bsky daha organik eğilimli, ancak zh/en izlenmeli",
    ]
    for i, note in enumerate(notes):
        add_text(slide, "• " + note, 9.95, 2.38 + i * 0.55, 2.7, 0.34, size=12, color=COLORS["muted"])
    add_footer(slide, 6)

    # 7. Themes and time
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Tema ve zaman: risk Crypto ve Politics etrafında yoğunlaşıyor", "Dashboard yalnız skor dağılımı değil, manipülasyon bağlamını da gösteriyor.")
    top_theme = theme.sort_values("mean_manipulative_score", ascending=False).head(6)
    chart_data = CategoryChartData()
    chart_data.categories = top_theme["primary_theme"].astype(str).tolist()
    chart_data.add_series("Mean score", top_theme["mean_manipulative_score"].tolist())
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.78), Inches(1.75), Inches(5.7), Inches(4.2), chart_data).chart
    style_chart(chart, 8)
    chart.has_legend = False
    chart.value_axis.maximum_scale = 0.75
    chart.plots[0].series[0].format.fill.solid()
    chart.plots[0].series[0].format.fill.fore_color.rgb = rgb(COLORS["risk"])
    trend_sorted = trend.sort_values("day")
    chart_data = CategoryChartData()
    chart_data.categories = trend_sorted["day"].astype(str).str.slice(5).tolist()
    chart_data.add_series("Manipulative rate", trend_sorted["manipulative_rate"].tolist())
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(7.0), Inches(1.75), Inches(5.2), Inches(3.05), chart_data).chart
    style_chart(chart, 8)
    chart.has_legend = False
    chart.value_axis.maximum_scale = max(0.25, trend_sorted["manipulative_rate"].max() + 0.03)
    chart.plots[0].series[0].format.line.color.rgb = rgb(COLORS["blue"])
    add_metric(slide, "68.1%", "Crypto manipülatif karar oranı", 7.15, 5.25, 2.4, COLORS["risk"])
    add_metric(slide, "50.4%", "Politics manipülatif karar oranı", 10.0, 5.25, 2.3, COLORS["blue"])
    add_footer(slide, 7)

    # 8. Author and duplicate networks
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Ağ örüntüleri: tekrar, yüksek hacim ve yüksek skor aynı yerde buluşuyor", "Author risk ve duplicate cluster listeleri manipülasyon ağını operasyonel olarak işaretler.")
    top_authors = author.head(5)
    add_text(slide, "Top author risk", 0.78, 1.65, 2.3, 0.28, size=15, bold=True, color=COLORS["risk"])
    x0, y0 = 0.78, 2.08
    headers = ["author_hash", "rows", "M rate", "risk"]
    widths = [3.2, 0.85, 0.85, 0.85]
    for j, htxt in enumerate(headers):
        add_text(slide, htxt, x0 + sum(widths[:j]), y0, widths[j], 0.2, size=8, bold=True, color=COLORS["muted"])
    for i, row in top_authors.iterrows():
        yy = y0 + 0.38 + i * 0.48
        fill = "FFFFFF" if i % 2 == 0 else "F0ECE5"
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0), Inches(yy - 0.08), Inches(sum(widths)), Inches(0.42))
        rect.fill.solid()
        rect.fill.fore_color.rgb = rgb(fill)
        rect.line.fill.background()
        vals = [str(row["author_hash"])[:10] + "...", f"{int(row['rows']):,}", pct(row["manipulative_rate"]), f"{row['author_risk_score']:.3f}"]
        for j, val in enumerate(vals):
            add_text(slide, val, x0 + sum(widths[:j]) + 0.05, yy, widths[j] - 0.1, 0.18, size=8, bold=(j == 3), color=(COLORS["risk"] if j == 3 else COLORS["ink"]))
    add_text(slide, "Duplicate cluster sinyali", 7.0, 1.65, 2.6, 0.28, size=15, bold=True, color=COLORS["trust"])
    for i, row in dup.head(4).iterrows():
        yy = 2.13 + i * 0.82
        add_panel(slide, 7.0, yy, 5.05, 0.62, "FFFFFF")
        text_excerpt = str(row.get("sample_text", "")).replace("\n", " ")[:118]
        count_value = row.get("rows", 0)
        mean_value = row.get("mean_score", 0.0)
        max_value = row.get("max_score", 0.0)
        count = int(count_value) if not pd.isna(count_value) else 0
        mean_score = float(mean_value) if not pd.isna(mean_value) else 0.0
        max_score = float(max_value) if not pd.isna(max_value) else 0.0
        add_text(slide, f"{count} tekrar | ort {mean_score:.3f} | max {max_score:.3f}", 7.18, yy + 0.09, 2.22, 0.18, size=7, bold=True, color=COLORS["risk"])
        add_text(slide, text_excerpt, 9.48, yy + 0.08, 2.33, 0.28, size=7, color=COLORS["muted"])
    add_text(slide, "Jüri mesajı: Sistem sadece tekil metni sınıflandırmıyor; hangi kümelerin organize risk ürettiğini de gösteriyor.", 1.05, 6.2, 10.9, 0.42, size=18, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 8)

    # 9. Live demo
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Canlı tahmin: gizli metin saniyeler içinde açıklamalı karar alır", "CLI ve Streamlit aynı model bundle’ını ve aynı threshold değerini kullanır.")
    add_text(slide, "Jüri metni", 0.9, 1.7, 1.4, 0.28, size=13, bold=True)
    add_panel(slide, 0.9, 2.12, 3.2, 1.25, "FFFFFF")
    add_text(slide, "“Free crypto giveaway, join our Telegram now.”", 1.08, 2.45, 2.84, 0.42, size=17, bold=True, color=COLORS["dark"])
    add_text(slide, "→", 4.38, 2.55, 0.35, 0.35, size=27, bold=True, color=COLORS["risk"])
    add_text(slide, "V11 inference", 4.92, 1.7, 1.6, 0.28, size=13, bold=True)
    add_panel(slide, 4.92, 2.12, 2.6, 1.25, "FFFFFF")
    add_text(slide, "tokenizer\nXLM-RoBERTa\nthreshold + review band", 5.18, 2.32, 2.05, 0.65, size=14, color=COLORS["muted"])
    add_text(slide, "→", 7.78, 2.55, 0.35, 0.35, size=27, bold=True, color=COLORS["risk"])
    add_text(slide, "Çıktı", 8.3, 1.7, 1.2, 0.28, size=13, bold=True)
    add_panel(slide, 8.3, 2.12, 3.7, 1.25, "FFFFFF")
    add_text(slide, "Karar: Manipülatif\nSkor: 0-1\nGerekçe + yakın manuel örnek", 8.55, 2.32, 3.15, 0.65, size=14, bold=True, color=COLORS["risk"])
    add_text(slide, "CLI", 1.0, 4.35, 0.55, 0.22, size=12, bold=True, color=COLORS["blue"])
    add_text(slide, "python version11_bert\\scripts\\v11_inference.py --text \"...\" --json", 1.0, 4.72, 8.7, 0.25, size=15, color=COLORS["dark"])
    add_text(slide, "Dashboard", 1.0, 5.38, 1.1, 0.22, size=12, bold=True, color=COLORS["trust"])
    add_text(slide, "Tek metin tahmini, batch CSV tahmini, platform/dil/tema/zaman/author risk haritası", 1.0, 5.75, 9.6, 0.28, size=15, color=COLORS["muted"])
    add_footer(slide, 9)

    # 10. Rubric mapping
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Puanlama karşılığı: model değil, karar destek sistemi", "Sunumda her rubric maddesi somut artefact ve canlı ekranla bağlanacak.")
    rows = [
        ("Analitik derinlik", "Platform/dil/tema/zaman/author/duplicate cluster analizleri", "15-20 hedefi"),
        ("Model + canlı tahmin", "External PR-AUC 0.892, M recall 0.899, CLI + dashboard inference", "29-40 hedefi"),
        ("Dashboard + sunum", "Risk haritası, author listesi, high-risk örnekler, batch tahmin", "15-20 hedefi"),
        ("Teknik uygulanabilirlik", "CUDA train, resumable 5M scoring, Drive artefact, GitHub code", "15-20 hedefi"),
    ]
    x0, y0 = 0.88, 1.78
    for i, (rubric, proof, target) in enumerate(rows):
        yy = y0 + i * 1.08
        add_text(slide, rubric, x0, yy, 2.3, 0.26, size=15, bold=True, color=COLORS["dark"])
        add_text(slide, proof, x0 + 2.8, yy, 6.3, 0.32, size=15, color=COLORS["muted"])
        add_chip(slide, target, x0 + 9.55, yy - 0.02, 1.7, COLORS["risk" if i < 3 else "trust"])
        add_rule(slide, x0, yy + 0.63, 10.9)
    add_text(slide, "Kapanış cümlesi", 0.95, 6.2, 1.6, 0.25, size=11, bold=True, color=COLORS["risk"])
    add_text(slide, "V11, tekil içerik kararını canlı verirken aynı skorları kullanarak manipülasyonun nerede, hangi dilde, hangi temada ve hangi author kümelerinde yoğunlaştığını gösterir.", 2.55, 6.08, 8.9, 0.5, size=17, bold=True)
    add_footer(slide, 10)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    OUT_PLAN.write_text(build_plan_text(summary, hold, ext), encoding="utf-8")
    validate_pptx(OUT_PPTX)
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {OUT_PLAN}")


def build_plan_text(summary: dict, holdout: dict, external: dict) -> str:
    return f"""# V11 Pitch Deck - Slayt Akışı

## 1. Kapak
Amaç: V11'in ne olduğunu tek cümlede vermek.
Bileşenler: 5M skorlanan içerik, external ROC-AUC, M recall, DataLeague final bağlamı.

## 2. İstenen Üç Çıktı
Amaç: README'deki üç teslimatı aynı sistemde karşıladığımızı göstermek.
Bileşenler: Organiklik skoru, manipülasyon haritası, canlı inference.

## 3. Mimari
Amaç: Canlı kararın text-only olduğunu, metadata'nın dashboard analitiğinde kullanıldığını açıklamak.
Bileşenler: 60k manuel M/O -> XLM-RoBERTa -> threshold -> 5M scoring -> dashboard.

## 4. Model Başarısı
Amaç: Modelin çalıştığını ve manipülatifi kaçırmama yaklaşımını kanıtlamak.
Bileşenler: Holdout/external ROC-AUC, PR-AUC, M recall, M F1; confusion matrix.
Holdout ROC-AUC: {holdout['roc_auc']:.4f}
External ROC-AUC: {external['roc_auc']:.4f}
External M recall: {external['recall_M']:.4f}

## 5. 5M Full Scoring
Amaç: Sadece örneklem değil, tüm veri üzerinde skor üretildiğini göstermek.
Bileşenler: {summary['rows_scored']:,} satır, manipülatif karar oranı {summary['decision_rates']['manipulative_rate']:.2%}, review oranı {summary['decision_rates']['review_rate']:.2%}.

## 6. Platform x Dil Haritası
Amaç: Dashboard puanını yükselten karar destek bileşenini göstermek.
Bileşenler: yüksek hacimli platform/dil heatmap; x.com/ar ve YouTube/en bulguları.

## 7. Tema ve Zaman
Amaç: Manipülasyonun konu ve zaman boyutunu göstermek.
Bileşenler: Cryptocurrency, Politics, Investing riskleri; günlük manipülatif oran trendi.

## 8. Author ve Duplicate Ağları
Amaç: Bot/spam/anomali örüntülerini yarışma problemiyle ilişkilendirmek.
Bileşenler: top risky authors, yüksek risk duplicate cluster örnekleri.

## 9. Canlı Demo
Amaç: Jürinin vereceği unseen metni nasıl çalıştıracağımızı göstermek.
Bileşenler: CLI komutu, Streamlit demo, skor, karar, review flag, gerekçeler.

## 10. Rubric Haritası
Amaç: Her puan kategorisini somut artefact ile bağlamak.
Bileşenler: analitik derinlik, model+canlı tahmin, dashboard+sunum, teknik uygulanabilirlik.
"""


def validate_pptx(path: Path) -> None:
    prs = Presentation(path)
    if len(prs.slides) != 10:
        raise RuntimeError(f"Expected 10 slides, got {len(prs.slides)}")
    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()
        if any("model.safetensors" in n or "v11_full_scored" in n for n in names):
            raise RuntimeError("Large artefact accidentally embedded in deck")
        slide_xml = [n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        if len(slide_xml) != 10:
            raise RuntimeError(f"Expected 10 slide XML files, got {len(slide_xml)}")


if __name__ == "__main__":
    build_deck()
